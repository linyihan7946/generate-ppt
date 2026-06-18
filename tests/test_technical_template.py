from pathlib import Path
from zipfile import ZipFile
import re
import xml.etree.ElementTree as ET

from pptx import Presentation

from generate_ppt.pipeline import generate_ppt


def test_technical_template_generates_landscape_deck_with_visible_animation_density(tmp_path: Path):
    source = tmp_path / "sample.txt"
    source.write_text(
        "\n".join(
            [
                "AI编程全流程",
                "用户输入需求",
                "LLM理解需求",
                "LLM读取文档",
                "AGENTS.md README.md",
                "文档作用",
                "提供一个项目的导航地图",
                "一方面提升速度",
                "一方面减少token消耗",
                "LLM搜索相关代码",
                "建立上下文（非常关键）",
                "搞清楚上下游关系、数据流、调用链和影响范围。",
                "制定计划",
                "循环修改代码，直到验证通过",
                "修改代码",
                "验证修改",
                "编译 类型检测 单元测试 集成测试",
                "测试失败 继续修改",
                "测试成功 退出循环",
            ]
        ),
        encoding="utf-8",
    )

    output = generate_ppt(source, "technical-no-image", tmp_path)
    prs = Presentation(output)

    assert round(prs.slide_width / 914400, 3) == 13.333
    assert round(prs.slide_height / 914400, 3) == 7.5
    assert len(prs.slides) >= 12

    all_text = "\n".join(shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text"))
    assert "技术文档讲解" in all_text
    assert "问题" in all_text
    assert "方案" in all_text
    assert "完整闭环" in all_text

    with ZipFile(output) as deck:
        slide_names = [name for name in deck.namelist() if re.match(r"ppt/slides/slide\d+\.xml$", name)]
        timing_slides = 0
        fade_effects = 0
        for name in slide_names:
            xml = deck.read(name)
            ET.fromstring(xml)
            text = xml.decode("utf-8", errors="ignore")
            timing_slides += "<p:timing>" in text
            fade_effects += text.count("<p:animEffect")

    assert timing_slides == len(prs.slides)
    assert fade_effects >= 80
