from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from generate_ppt.slide_model import DeckSpec, SlideKind, SlideSpec


W = Inches(13.333)
H = Inches(7.5)

COLORS = {
    "bg": RGBColor(8, 12, 18),
    "panel": RGBColor(18, 25, 34),
    "panel2": RGBColor(25, 34, 45),
    "grid": RGBColor(36, 55, 68),
    "text": RGBColor(238, 244, 250),
    "muted": RGBColor(150, 165, 178),
    "blue": RGBColor(65, 165, 255),
    "green": RGBColor(48, 210, 125),
    "red": RGBColor(255, 78, 91),
    "yellow": RGBColor(255, 210, 75),
    "cyan": RGBColor(61, 228, 235),
    "purple": RGBColor(159, 114, 255),
    "orange": RGBColor(255, 144, 72),
}


class PptRenderer:
    def render(self, deck: DeckSpec, output_path: Path) -> Path:
        prs = Presentation()
        prs.slide_width = W
        prs.slide_height = H
        for idx, spec in enumerate(deck.slides, 1):
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            self._background(slide)
            self._render_slide(slide, idx, spec)
            self._footer(slide, idx)
        output_path.parent.mkdir(parents=True, exist_ok=True)
        prs.save(output_path)
        return output_path

    def _render_slide(self, slide, idx: int, spec: SlideSpec) -> None:
        if spec.kind == SlideKind.COVER:
            self._cover(slide, spec)
        elif spec.kind == SlideKind.SECTION:
            self._section(slide, spec)
        elif spec.kind == SlideKind.FLOW:
            self._flow(slide, spec)
        elif spec.kind == SlideKind.COMPARE:
            self._compare(slide, spec)
        elif spec.kind == SlideKind.PYRAMID:
            self._pyramid(slide, spec)
        elif spec.kind == SlideKind.RELATION:
            self._relation(slide, spec)
        elif spec.kind == SlideKind.PROBLEM_CHAIN:
            self._problem_chain(slide, spec)
        elif spec.kind == SlideKind.SCREENSHOT:
            self._screenshot(slide, spec)
        elif spec.kind == SlideKind.SUMMARY:
            self._summary(slide, spec)
        else:
            self._bullets(slide, spec)

    def _background(self, slide) -> None:
        rect = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, W, H)
        rect.fill.solid()
        rect.fill.fore_color.rgb = COLORS["bg"]
        rect.line.fill.background()

        step = Inches(0.62)
        x = 0
        while x < W:
            line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x, 0, x, H)
            line.line.color.rgb = COLORS["grid"]
            line.line.transparency = 76
            line.line.width = Pt(0.4)
            x += step
        y = 0
        while y < H:
            line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, 0, y, W, y)
            line.line.color.rgb = COLORS["grid"]
            line.line.transparency = 80
            line.line.width = Pt(0.4)
            y += step

    def _cover(self, slide, spec: SlideSpec) -> None:
        self._badge(slide, "TECH DECK", Inches(0.72), Inches(0.58), Inches(1.75), Inches(0.36), "blue", size=13, name="anim_01_badge")
        self._text(slide, spec.title[:28], Inches(0.75), Inches(1.38), Inches(6.2), Inches(1.45), 47, bold=True, name="anim_02_title")
        subtitle = spec.subtitle or "技术文档讲解"
        self._text(slide, subtitle[:36], Inches(0.78), Inches(3.0), Inches(5.6), Inches(0.72), 28, "cyan", bold=True, name="anim_03_subtitle")
        self._chip(slide, Inches(8.4), Inches(1.8), Inches(3.0), Inches(1.55), "anim_04_")
        labels = spec.items[:3] or ["理解", "拆解", "验证"]
        xs = [Inches(7.4), Inches(9.2), Inches(11.0)]
        for i, label in enumerate(labels):
            self._card(slide, label[:6], xs[i], Inches(4.65), Inches(1.38), Inches(0.58), ["blue", "green", "yellow"][i], f"anim_1{i}_cover")

    def _section(self, slide, spec: SlideSpec) -> None:
        self._badge(slide, spec.subtitle or "重点", Inches(0.75), Inches(0.62), Inches(2.2), Inches(0.52), spec.accent, size=20, name="anim_01_badge")
        self._text(slide, spec.title, Inches(0.75), Inches(1.65), Inches(5.5), Inches(1.05), 46, bold=True, name="anim_02_title")
        if spec.highlight:
            self._text(slide, spec.highlight, Inches(0.85), Inches(3.1), Inches(5.7), Inches(0.95), 30, "cyan", bold=True, align=PP_ALIGN.CENTER, name="anim_03_highlight")
        for i, item in enumerate(spec.items[:3]):
            self._outline(slide, item, Inches(7.0), Inches(1.35 + i * 1.3), Inches(4.9), Inches(0.78), ["blue", "yellow", "green"][i % 3], f"anim_{i+4:02d}_item")

    def _bullets(self, slide, spec: SlideSpec) -> None:
        self._title(slide, spec)
        for i, item in enumerate(spec.items[:5]):
            row = i // 2
            col = i % 2
            x = Inches(0.95 + col * 6.0)
            y = Inches(2.1 + row * 1.28)
            self._badge(slide, f"{i + 1}", x, y, Inches(0.52), Inches(0.52), spec.accent, size=19, name=f"anim_{i+1:02d}_num")
            self._text(slide, item, x + Inches(0.72), y - Inches(0.05), Inches(4.85), Inches(0.65), 24, bold=True, name=f"anim_{i+1:02d}_item")

    def _flow(self, slide, spec: SlideSpec) -> None:
        self._title(slide, spec)
        items = spec.items[:7]
        start_x = Inches(0.65)
        y = Inches(3.15)
        for i, item in enumerate(items):
            x = start_x + i * Inches(1.78)
            color = ["blue", "cyan", "yellow", "purple", "orange", "green", "red"][i % 7]
            self._card(slide, item, x, y, Inches(1.45), Inches(0.82), color, f"anim_{i+1:02d}_node")
            if i < len(items) - 1:
                self._line(slide, x + Inches(1.5), y + Inches(0.41), x + Inches(1.75), y + Inches(0.41), color, f"anim_{i+1:02d}_arrow")

    def _compare(self, slide, spec: SlideSpec) -> None:
        self._title(slide, spec)
        left = spec.items[0] if spec.items else "问题"
        right = spec.items[1] if len(spec.items) > 1 else "方案"
        self._card(slide, left, Inches(1.0), Inches(2.65), Inches(4.0), Inches(1.1), "red", "anim_01_left")
        self._line(slide, Inches(5.25), Inches(3.2), Inches(6.25), Inches(3.2), "blue", "anim_02_arrow")
        self._card(slide, right, Inches(6.55), Inches(2.65), Inches(4.0), Inches(1.1), "green", "anim_03_right")
        for i, item in enumerate(spec.items[2:4]):
            self._badge(slide, item, Inches(2.0 + i * 5.1), Inches(4.7), Inches(3.0), Inches(0.62), ["yellow", "cyan"][i], size=23, name=f"anim_0{i+4}_badge")
        if spec.highlight:
            self._text(slide, spec.highlight, Inches(1.2), Inches(5.82), Inches(10.8), Inches(0.62), 29, "yellow", bold=True, align=PP_ALIGN.CENTER, name="anim_06_conclusion")

    def _pyramid(self, slide, spec: SlideSpec) -> None:
        self._title(slide, spec)
        items = spec.items[:5]
        widths = [9.2, 7.6, 6.0, 4.4, 3.0]
        for i, item in enumerate(items):
            w = Inches(widths[min(i, len(widths) - 1)])
            x = (W - w) / 2
            y = Inches(1.95 + i * 0.85)
            self._card(slide, item, x, y, w, Inches(0.66), ["blue", "cyan", "yellow", "orange", "green"][i % 5], f"anim_{i+1:02d}_level")

    def _relation(self, slide, spec: SlideSpec) -> None:
        self._title(slide, spec)
        center = spec.items[0] if spec.items else spec.title
        self._card(slide, center, Inches(5.35), Inches(3.2), Inches(2.6), Inches(0.9), "cyan", "anim_01_center")
        labels = (spec.items[1:5] + ["上游输入", "下游影响", "数据流", "调用链"])[:4]
        positions = [(1.1, 2.1, "blue"), (9.9, 2.1, "red"), (1.1, 5.1, "yellow"), (9.9, 5.1, "purple")]
        for i, (label, (x, y, color)) in enumerate(zip(labels, positions), 2):
            self._outline(slide, label, Inches(x), Inches(y), Inches(2.25), Inches(0.72), color, f"anim_{i:02d}_rel")
            self._line(slide, Inches(x + 1.12), Inches(y + 0.72), Inches(6.65), Inches(3.65), color, f"anim_{i:02d}_line")

    def _problem_chain(self, slide, spec: SlideSpec) -> None:
        self._title(slide, spec)
        labels = ["问题", "原因", "方案", "结果"]
        colors = ["red", "orange", "green", "yellow"]
        items = (spec.items + ["待补充"] * 4)[:4]
        for i, (label, item) in enumerate(zip(labels, items)):
            x = Inches(0.75 + i * 3.1)
            y = Inches(3.0)
            self._badge(slide, label, x, y - Inches(0.62), Inches(1.1), Inches(0.48), colors[i], size=18, name=f"anim_{i+1:02d}_label")
            self._card(slide, item, x, y, Inches(2.45), Inches(0.78), colors[i], f"anim_{i+1:02d}_item")
            if i < 3:
                self._line(slide, x + Inches(2.5), y + Inches(0.39), x + Inches(2.95), y + Inches(0.39), colors[i], f"anim_{i+1:02d}_arrow")

    def _screenshot(self, slide, spec: SlideSpec) -> None:
        self._title(slide, spec)
        self._screenshot_box(slide, spec.screenshot_hint or "预留：软件界面 / 配置页面 / 操作步骤截图", Inches(0.8), Inches(1.9), Inches(7.5), Inches(4.45), "anim_01_screenshot")
        for i, item in enumerate(spec.items[:2]):
            self._badge(slide, item, Inches(8.75), Inches(2.45 + i * 1.15), Inches(3.0), Inches(0.62), ["green", "yellow"][i], size=22, name=f"anim_0{i+2}_tip")

    def _summary(self, slide, spec: SlideSpec) -> None:
        self._text(slide, spec.title, Inches(0.75), Inches(0.95), Inches(11.85), Inches(0.95), 45, bold=True, align=PP_ALIGN.CENTER, name="anim_01_title")
        for i, item in enumerate(spec.items[:7]):
            x = Inches(0.95 + i * 1.72)
            y = Inches(3.25)
            self._badge(slide, item, x, y, Inches(1.28), Inches(0.54), ["blue", "yellow", "purple", "orange", "cyan", "green", "red"][i % 7], size=18, name=f"anim_{i+2:02d}_summary")
            if i < len(spec.items[:7]) - 1:
                self._line(slide, x + Inches(1.33), y + Inches(0.27), x + Inches(1.58), y + Inches(0.27), "cyan", f"anim_{i+2:02d}_arrow")
        if spec.highlight:
            self._text(slide, spec.highlight, Inches(1.2), Inches(5.25), Inches(10.8), Inches(0.8), 31, "yellow", bold=True, align=PP_ALIGN.CENTER, name="anim_20_end")

    def _title(self, slide, spec: SlideSpec) -> None:
        if spec.subtitle:
            self._text(slide, spec.subtitle, Inches(0.65), Inches(0.42), Inches(11.9), Inches(0.35), 15, spec.accent, bold=True, name="anim_00_kicker")
        self._text(slide, spec.title, Inches(0.65), Inches(0.78), Inches(11.9), Inches(0.72), 33, bold=True, name="anim_00_title")

    def _footer(self, slide, idx: int) -> None:
        self._text(slide, f"{idx:02d}", Inches(12.55), Inches(6.92), Inches(0.42), Inches(0.25), 10, "muted", bold=True, align=PP_ALIGN.RIGHT)

    def _text(self, slide, text: str, x, y, w, h, size=28, color="text", bold=False, align=PP_ALIGN.LEFT, name: str | None = None):
        box = slide.shapes.add_textbox(x, y, w, h)
        tf = box.text_frame
        tf.clear()
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = 0
        tf.margin_right = 0
        tf.margin_top = 0
        tf.margin_bottom = 0
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.name = "Microsoft YaHei UI"
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = COLORS[color]
        if name:
            box.name = name
        return box

    def _badge(self, slide, text, x, y, w, h, color, size=17, name=None):
        shp = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h)
        shp.fill.solid()
        shp.fill.fore_color.rgb = COLORS[color]
        shp.line.fill.background()
        self._shape_text(shp, text, size, "bg", True)
        if name:
            shp.name = name
        return shp

    def _card(self, slide, text, x, y, w, h, color, name=None):
        shp = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h)
        shp.fill.solid()
        shp.fill.fore_color.rgb = COLORS["panel"]
        shp.line.color.rgb = COLORS[color]
        shp.line.width = Pt(1.7)
        self._shape_text(shp, text, 24, "text", True)
        if name:
            shp.name = name
        return shp

    def _outline(self, slide, text, x, y, w, h, color, name=None):
        shp = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h)
        shp.fill.solid()
        shp.fill.fore_color.rgb = COLORS["bg"]
        shp.line.color.rgb = COLORS[color]
        shp.line.width = Pt(2.1)
        self._shape_text(shp, text, 22, "text", True)
        if name:
            shp.name = name
        return shp

    def _shape_text(self, shp, text, size, color, bold) -> None:
        tf = shp.text_frame
        tf.clear()
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = Inches(0.08)
        tf.margin_right = Inches(0.08)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = text
        run.font.name = "Microsoft YaHei UI"
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = COLORS[color]

    def _line(self, slide, x1, y1, x2, y2, color, name=None):
        line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
        line.line.color.rgb = COLORS[color]
        line.line.width = Pt(2.2)
        line.line.end_arrowhead = True
        if name:
            line.name = name
        return line

    def _chip(self, slide, x, y, w, h, name_prefix=""):
        self._outline(slide, "AI", x, y, w, h, "blue", f"{name_prefix}chip")
        for i in range(4):
            yy = y + Inches(0.18) + i * (h - Inches(0.36)) / 3
            self._line(slide, x - Inches(0.18), yy, x, yy, "cyan", f"{name_prefix}pin_l{i}").line.end_arrowhead = False
            self._line(slide, x + w, yy, x + w + Inches(0.18), yy, "cyan", f"{name_prefix}pin_r{i}").line.end_arrowhead = False

    def _screenshot_box(self, slide, label, x, y, w, h, name=None):
        frame = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h)
        frame.fill.solid()
        frame.fill.fore_color.rgb = RGBColor(13, 18, 25)
        frame.line.color.rgb = COLORS["cyan"]
        frame.line.width = Pt(2.2)
        bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, x, y, w, Inches(0.38))
        bar.fill.solid()
        bar.fill.fore_color.rgb = COLORS["panel2"]
        bar.line.fill.background()
        self._shape_text(frame, label, 24, "muted", True)
        if name:
            frame.name = name
        return frame
